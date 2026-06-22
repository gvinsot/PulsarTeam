"""PowerPoint (.pptx) — read, edit-in-place, generate. Backed by python-pptx.

Editing is round-trip: existing slides/shapes are preserved; the LLM targets a
slide by index and replaces text or appends a slide. Reads expose per-slide
title, body text and speaker notes, plus a markdown rendering.
"""
from __future__ import annotations

from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from ..util import ToolError, check_size, default_output, resolve_output, resolve_path


def _slide_title(slide) -> str | None:
    if slide.shapes.title is not None:
        return slide.shapes.title.text
    return None


def _slide_texts(slide) -> list[str]:
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs) or para.text
                if text.strip():
                    out.append(text)
    return out


def _notes(slide) -> str | None:
    if slide.has_notes_slide:
        return slide.notes_slide.notes_text_frame.text or None
    return None


def read_pptx(path: str) -> dict[str, Any]:
    src = resolve_path(path)
    check_size(src)
    prs = Presentation(str(src))
    slides = []
    md = []
    for i, slide in enumerate(prs.slides):
        title = _slide_title(slide)
        body = _slide_texts(slide)
        notes = _notes(slide)
        slides.append({"index": i, "title": title, "body": body, "notes": notes})
        md.append(f"## Slide {i + 1}: {title or '(no title)'}")
        md.extend(f"- {b}" for b in body)
        if notes:
            md.append(f"> notes: {notes}")
    return {
        "format": "pptx",
        "path": str(src),
        "slide_count": len(slides),
        "slides": slides,
        "markdown": "\n\n".join(md),
    }


def get_outline(path: str) -> dict[str, Any]:
    src = resolve_path(path)
    prs = Presentation(str(src))
    slides = [{"index": i, "title": _slide_title(s)} for i, s in enumerate(prs.slides)]
    return {"format": "pptx", "path": str(src), "slides": slides}


def _replace_in_slide(slide, find: str, repl: str) -> int:
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            joined = "".join(run.text for run in para.runs)
            if find in joined and para.runs:
                para.runs[0].text = joined.replace(find, repl)
                for r in para.runs[1:]:
                    r.text = ""
                count += 1
    return count


def edit_pptx(path: str, operations: list[dict[str, Any]], output_path: str | None = None) -> dict[str, Any]:
    """Operations:
    - {op:"replace_text", find, replace, slide?}   slide index optional (all if omitted)
    - {op:"set_title", slide, text}
    - {op:"set_notes", slide, text}
    - {op:"add_slide", layout?:int, title?, bullets?:[...], notes?}
    """
    src = resolve_path(path)
    if not isinstance(operations, list) or not operations:
        raise ToolError("operations must be a non-empty list")
    prs = Presentation(str(src))
    slides = list(prs.slides)
    applied = []
    for op in operations:
        kind = op.get("op")
        if kind == "replace_text":
            find = op.get("find")
            if not find:
                raise ToolError("replace_text requires 'find'")
            repl = op.get("replace", "")
            targets = [slides[op["slide"]]] if "slide" in op else slides
            n = sum(_replace_in_slide(s, find, repl) for s in targets)
            applied.append(f"replace_text: {n} run-group(s)")
        elif kind == "set_title":
            s = slides[op["slide"]]
            if s.shapes.title is None:
                raise ToolError(f"slide {op['slide']} has no title placeholder")
            s.shapes.title.text = op.get("text", "")
            applied.append(f"set_title[{op['slide']}]")
        elif kind == "set_notes":
            s = slides[op["slide"]]
            s.notes_slide.notes_text_frame.text = op.get("text", "")
            applied.append(f"set_notes[{op['slide']}]")
        elif kind == "add_slide":
            _add_slide(prs, op)
            applied.append("add_slide")
        else:
            raise ToolError(f"unknown pptx op: {kind}")
    out = resolve_output(output_path) if output_path else default_output(src)
    prs.save(str(out))
    return {"ok": True, "applied": applied, "output_path": str(out)}


def _add_slide(prs, spec: dict[str, Any]) -> None:
    layout_idx = int(spec.get("layout", 1))
    layout_idx = min(layout_idx, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    if spec.get("title") and slide.shapes.title is not None:
        slide.shapes.title.text = str(spec["title"])
    bullets = spec.get("bullets") or []
    if bullets:
        body = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:  # not the title
                body = ph
                break
        if body is not None and body.has_text_frame:
            tf = body.text_frame
            tf.text = str(bullets[0])
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = str(b)
    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = str(spec["notes"])


def generate_pptx(output_path: str, spec: dict[str, Any]) -> dict[str, Any]:
    """Create a new deck.

    spec = {"slides": [{"layout?":int, "title": "...", "bullets": [...], "notes?": "..."}]}
    """
    out = resolve_output(output_path)
    prs = Presentation()
    slides = spec.get("slides")
    if not slides:
        raise ToolError("spec needs a non-empty 'slides' list")
    for s in slides:
        _add_slide(prs, s)
    prs.save(str(out))
    return {"ok": True, "output_path": str(out), "slide_count": len(slides)}
